
#!/usr/bin/env python3
"""
DEF Report Generator
Template-driven offensive tendency report generator.

Usage:
    python def_report_generator.py \
      --template "MASTER Offensive Breakdown Template.pptx" \
      --output "Opponent_Offensive_Breakdown.pptx" \
      --opponent "Opponent" \
      --excel "Game1.xlsx" "Game2.xlsx" "Game3.xlsx"

Core behavior:
- Reads one or more Hudl Excel/CSV exports.
- Filters to ODK = O by default.
- Normalizes columns.
- Builds consistent analysis objects.
- Opens the MASTER Offensive Breakdown Template.
- Populates existing slides and duplicates Formation/Situation slides.
- Preserves the template's theme, slide size, colors, table styling, and layouts as much as python-pptx allows.

Dependencies:
    pip install pandas openpyxl python-pptx
"""
from __future__ import annotations

import argparse
import copy
import math
import re
from pathlib import Path
from typing import Dict, List, Tuple, Any, Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


# -----------------------------
# CONFIG
# -----------------------------

DND_ORDER = [
    "1st & 10", "1st & 11+", "1st & <10",
    "2nd & 8+", "2nd & 4-7", "2nd & 1-3",
    "3rd & 8+", "3rd & 4-7", "3rd & 1-3",
    "4th & 8+", "4th & 4-7", "4th & 1-3"
]

FIELD_ZONE_ORDER = [
    "Coming Out", "Backed Up", "3 Down Territory", "Midfield",
    "4 Down Territory", "High Red Zone", "Low Red Zone", "Goal Line"
]

REQUIRED_CONCEPTS = {
    "ODK": ["ODK"],
    "DN": ["DN", "DOWN"],
    "DIST": ["DIST", "DISTANCE"],
    "YARD_LN": ["YARD LN", "YARD LINE", "BALL ON"],
    "HASH": ["HASH"],
    "OFF_STR": ["OFF STR", "OFF STRENGTH"],
    "PLAY_TYPE": ["PLAY TYPE", "PLAYTYPE"],
    "OFF_FORM": ["OFF FORM", "FORMATION"],
    "FORM_TAG": ["FORM TAG"],
    "MOTION": ["MOTION"],
    "MOTION_DIR": ["MOTION DIR"],
    "PERSONNEL": ["PERSONNEL", "PERS"],
    "BACKFIELD": ["BACKFIELD", "BACKFIELD SET"],
    "OFF_PLAY": ["OFF PLAY", "PLAY"],
    "PLAY_TAG": ["PLAY TAG"],
    "BOUNDARY_ROUTE": ["BOUNDARY ROUTE", "BNDY ROUTE"],
    "FIELD_ROUTE": ["FIELD ROUTE"],
    "EFF": ["EFF", "EFFICIENT"],
    "GNLS": ["GN/LS", "GAIN/LOSS", "YARDS", "GAIN"],
    "PLAY_DIR": ["PLAY DIR", "PLAY DIRECTION"],
    "PLAYER": ["PLAYER", "BALL CARRIER", "TARGET", "TOUCH"],
    "P10": ["P&10", "P & 10", "P-10", "P – 10"],
}


# -----------------------------
# UTILS
# -----------------------------

def _fmt_count(x: float) -> str:
    """Format count-like values without trailing .0."""
    try:
        if pd.isna(x):
            return "0"
        xf = float(x)
        if xf.is_integer():
            return str(int(xf))
        return f"{xf:.1f}"
    except Exception:
        return str(x)

def pct(n: float, d: float, whole: bool = True, show_counts: bool = True) -> str:
    """
    Percentage display standard.

    Every report percentage should show the numerator and denominator that created it:
        25% (2/8)

    This makes the PowerPoint more useful for coaches and athletes because it shows
    whether the percentage is based on a strong or limited sample.
    """
    if d is None or d == 0 or pd.isna(d):
        return "N/A"
    n = 0 if n is None or pd.isna(n) else n
    val = 100.0 * float(n) / float(d)
    base = f"{val:.0f}%" if whole else f"{val:.1f}%"
    if show_counts:
        return f"{base} ({_fmt_count(n)}/{_fmt_count(d)})"
    return base

def pct_value(text: Any) -> float:
    """Extract the leading numeric percentage from strings like '75% (6/8)'."""
    m = re.search(r"-?\d+(?:\.\d+)?", str(text))
    return float(m.group(0)) if m else -999.0

def ypp(yards: float, calls: float) -> str:
    if calls is None or calls == 0 or pd.isna(calls):
        return "N/A"
    return f"{yards / calls:.1f}"

def clean_text(x: Any, blank: str = "-") -> str:
    if pd.isna(x):
        return blank
    s = str(x).strip()
    if not s:
        return blank
    # remove .0 on numeric labels like player/personnel if appropriate
    if re.fullmatch(r"\d+\.0", s):
        return s[:-2]
    return re.sub(r"\s+", " ", s).upper()

def clean_display(x: Any, blank: str = "-") -> str:
    s = clean_text(x, blank)
    return s

def to_num(x: Any) -> float:
    try:
        if pd.isna(x):
            return math.nan
        return float(x)
    except Exception:
        m = re.search(r"-?\d+(\.\d+)?", str(x))
        return float(m.group(0)) if m else math.nan

def norm_play_type(x: Any) -> str:
    s = clean_text(x)
    if s.startswith("R"):
        return "Run"
    if s.startswith("P"):
        return "Pass"
    return "Unknown"

def is_eff(x: Any) -> bool:
    s = clean_text(x, blank="")
    return s in {"Y", "YES", "1", "TRUE", "T", "EFFICIENT"}

def get_col(df: pd.DataFrame, aliases: List[str]) -> Optional[str]:
    upper_map = {str(c).strip().upper(): c for c in df.columns}
    for a in aliases:
        if a.strip().upper() in upper_map:
            return upper_map[a.strip().upper()]
    return None

def rename_columns(df: pd.DataFrame) -> pd.DataFrame:
    mapping = {}
    for concept, aliases in REQUIRED_CONCEPTS.items():
        col = get_col(df, aliases)
        if col is not None:
            mapping[col] = concept
    return df.rename(columns=mapping)

def read_files(paths: List[str]) -> pd.DataFrame:
    frames = []
    for p in paths:
        path = Path(p)
        if path.suffix.lower() == ".csv":
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
        df = rename_columns(df)
        df["SOURCE_FILE"] = path.name
        frames.append(df)
    if not frames:
        raise ValueError("No input files provided.")
    return pd.concat(frames, ignore_index=True, sort=False)

def field_zone(yard: Any) -> str:
    y = to_num(yard)
    if pd.isna(y):
        return "Unknown"
    # Knowledge Base definitions
    if -10 <= y <= -1:
        return "Coming Out"
    if -20 <= y <= -11:
        return "Backed Up"
    if -49 <= y <= -21:
        return "3 Down Territory"
    if y in {0, 50, -50}:
        return "Midfield"
    if 21 <= y <= 49:
        return "4 Down Territory"
    if 11 <= y <= 20:
        return "High Red Zone"
    if 5 <= y <= 10:
        return "Low Red Zone"
    if 1 <= y <= 5:
        return "Goal Line"
    return "Unknown"

def dnd_category(dn: Any, dist: Any) -> str:
    d = to_num(dn)
    y = to_num(dist)
    if pd.isna(d) or pd.isna(y):
        return "Unknown"
    d = int(d)
    if d == 1:
        if y == 10:
            return "1st & 10"
        if y > 10:
            return "1st & 11+"
        return "1st & <10"
    if d in {2,3,4}:
        prefix = {2:"2nd",3:"3rd",4:"4th"}[d]
        if y >= 8:
            return f"{prefix} & 8+"
        if 4 <= y <= 7:
            return f"{prefix} & 4-7"
        if 1 <= y <= 3:
            return f"{prefix} & 1-3"
    return "Unknown"

def add_features(df: pd.DataFrame, odk: str = "O") -> pd.DataFrame:
    if "ODK" in df.columns:
        df = df[df["ODK"].astype(str).str.strip().str.upper() == odk.upper()].copy()
    else:
        raise ValueError("ODK column is required for this offensive report.")

    for c in ["GNLS", "DN", "DIST", "YARD_LN"]:
        if c in df.columns:
            df[c] = df[c].apply(to_num)

    for c in ["OFF_FORM","FORM_TAG","MOTION","MOTION_DIR","PERSONNEL","BACKFIELD","OFF_PLAY","PLAY_TAG",
              "BOUNDARY_ROUTE","FIELD_ROUTE","HASH","OFF_STR","PLAY_DIR","PLAYER"]:
        if c not in df.columns:
            df[c] = "-"
        df[c] = df[c].apply(clean_display)

    df["PLAY_TYPE_NORM"] = df["PLAY_TYPE"].apply(norm_play_type) if "PLAY_TYPE" in df.columns else "Unknown"
    df["YARDS"] = df["GNLS"] if "GNLS" in df.columns else 0
    df["YARDS"] = df["YARDS"].fillna(0)
    df["EFF_BOOL"] = df["EFF"].apply(is_eff) if "EFF" in df.columns else False
    df["DND"] = df.apply(lambda r: dnd_category(r.get("DN"), r.get("DIST")), axis=1)
    df["FIELD_ZONE"] = df["YARD_LN"].apply(field_zone) if "YARD_LN" in df.columns else "Unknown"
    if "P10" in df.columns:
        df["P10_YN"] = df["P10"].apply(lambda x: "YES" if clean_text(x, "") in {"Y","YES","1","TRUE"} else "NO")
    else:
        df["P10_YN"] = df.apply(lambda r: "YES" if r.get("DND") == "1st & 10" else "NO", axis=1)

    # Keep only run/pass rows for denominator.
    df = df[df["PLAY_TYPE_NORM"].isin(["Run","Pass"])].copy()
    df["IS_RUN"] = df["PLAY_TYPE_NORM"] == "Run"
    df["IS_PASS"] = df["PLAY_TYPE_NORM"] == "Pass"
    df["IS_EXPLOSIVE"] = ((df["IS_RUN"] & (df["YARDS"] > 10)) | (df["IS_PASS"] & (df["YARDS"] > 15)))
    df["PLAYCALL_RUN"] = df["OFF_FORM"] + " / " + df["MOTION"] + " / " + df["OFF_PLAY"]
    df["PLAYCALL_PASS"] = df["OFF_FORM"] + " / " + df["MOTION"] + " / " + df["OFF_PLAY"] + " / " + df["FIELD_ROUTE"] + " / " + df["BOUNDARY_ROUTE"] + " / " + df["PLAY_TAG"]
    df["COMBO"] = df["OFF_FORM"] + " / " + df["MOTION"] + " / " + df["OFF_PLAY"]
    return df

def validate(df_raw: pd.DataFrame) -> List[str]:
    issues = []
    for needed in ["ODK","PLAY_TYPE","OFF_FORM","OFF_PLAY","GNLS"]:
        if needed not in df_raw.columns:
            issues.append(f"Missing expected column concept: {needed}")
    return issues


# -----------------------------
# ANALYTICS
# -----------------------------

def count_run_pass(g: pd.DataFrame) -> Tuple[int,int,int]:
    total = len(g)
    run = int(g["IS_RUN"].sum()) if total else 0
    pas = int(g["IS_PASS"].sum()) if total else 0
    return total, run, pas

def core_metrics(g: pd.DataFrame) -> Dict[str, Any]:
    total, run, pas = count_run_pass(g)
    return {
        "total": total,
        "run": run,
        "pass": pas,
        "run_pct": pct(run, total),
        "pass_pct": pct(pas, total),
        "ypp": ypp(g["YARDS"].sum(), total),
        "eff_pct": pct(g["EFF_BOOL"].sum(), total),
        "explosive": int(g["IS_EXPLOSIVE"].sum()),
        "explosive_pct": pct(g["IS_EXPLOSIVE"].sum(), total),
    }

def top_counts(df: pd.DataFrame, col: str, n: int = 3, denom: Optional[int] = None, play_type: Optional[str] = None) -> List[List[str]]:
    g = df
    if play_type == "Run":
        g = g[g["IS_RUN"]]
    elif play_type == "Pass":
        g = g[g["IS_PASS"]]
    if g.empty or col not in g.columns:
        return []
    denom = denom or len(df)
    out = []
    for val, sub in g.groupby(col, dropna=False):
        calls = len(sub)
        out.append([str(val), calls, pct(calls, denom), ypp(sub["YARDS"].sum(), calls), sub["YARDS"].sum()])
    out.sort(key=lambda r: (-r[1], -float(r[3]) if r[3]!="N/A" else -999, r[0]))
    return [[r[0], str(r[1]), r[2], r[3]] for r in out[:n]]

def run_pass_table(df: pd.DataFrame, col: str, order: Optional[List[str]]=None, max_rows: int=99) -> List[List[str]]:
    rows = []
    keys = order if order else list(df[col].dropna().unique())
    for k in keys:
        sub = df[df[col] == k]
        if len(sub) == 0:
            continue
        m = core_metrics(sub)
        rows.append([str(k), str(m["total"]), str(m["run"]), m["run_pct"], str(m["pass"]), m["pass_pct"], m["ypp"]])
    if not order:
        rows.sort(key=lambda r: (-int(r[1]), r[0]))
    return rows[:max_rows]

def dnd_table(df: pd.DataFrame, max_rows: int=99) -> List[List[str]]:
    rows = []
    for k in DND_ORDER:
        sub = df[df["DND"] == k]
        if sub.empty: continue
        m = core_metrics(sub)
        rows.append([k, str(m["total"]), m["run_pct"], m["pass_pct"], m["ypp"], m["eff_pct"]])
    return rows[:max_rows]

def field_zone_table(df: pd.DataFrame, max_rows: int=99) -> List[List[str]]:
    rows=[]
    for k in FIELD_ZONE_ORDER:
        sub = df[df["FIELD_ZONE"] == k]
        if sub.empty: continue
        m=core_metrics(sub)
        rows.append([k, str(m["total"]), m["run_pct"], m["pass_pct"], m["ypp"], m["eff_pct"]])
    return rows[:max_rows]

def p10_table(df: pd.DataFrame) -> List[List[str]]:
    rows=[]
    for k in ["NO","YES"]:
        sub=df[df["P10_YN"]==k]
        if sub.empty: continue
        m=core_metrics(sub)
        rows.append([k, str(m["total"]), str(m["run"]), m["run_pct"], str(m["pass"]), m["pass_pct"], m["ypp"]])
    return rows

def formation_summary(df: pd.DataFrame) -> List[List[str]]:
    rows=[]
    for form, sub in df.groupby("OFF_FORM"):
        m=core_metrics(sub)
        top_play = sub["OFF_PLAY"].value_counts().index[0] if len(sub) else "-"
        rows.append([form, str(m["total"]), str(m["run"]), m["run_pct"], str(m["pass"]), m["pass_pct"], top_play])
    rows.sort(key=lambda r: (-int(r[1]), r[0]))
    return rows

def combo_summary(df: pd.DataFrame, n=10) -> List[List[str]]:
    vc = df["COMBO"].value_counts()
    return [[idx, str(cnt), pct(cnt, len(df))] for idx,cnt in vc.head(n).items()]

def formation_detail(df: pd.DataFrame, form: str) -> Dict[str, Any]:
    sub=df[df["OFF_FORM"]==form]
    m=core_metrics(sub)
    # alert = most common OFF_PLAY
    alert = sub["OFF_PLAY"].value_counts().index[0] if len(sub) else "-"
    alert_pct = pct((sub["OFF_PLAY"]==alert).sum(), len(sub)) if len(sub) else "N/A"
    return {
        "formation": form,
        "metrics": m,
        "run_plays": top_counts(sub, "OFF_PLAY", 3, denom=len(sub), play_type="Run"),
        "field_routes": top_counts(sub[sub["IS_PASS"]], "FIELD_ROUTE", 3, denom=max(1, len(sub[sub["IS_PASS"]]))) if len(sub[sub["IS_PASS"]]) else [],
        "boundary_routes": top_counts(sub[sub["IS_PASS"]], "BOUNDARY_ROUTE", 3, denom=max(1, len(sub[sub["IS_PASS"]]))) if len(sub[sub["IS_PASS"]]) else [],
        "motions": motion_or_backfield_rows(sub, "MOTION", max_items=3),
        "backfields": motion_or_backfield_rows(sub, "BACKFIELD", max_items=3),
        "alert": alert,
        "alert_pct": alert_pct,
    }

def motion_or_backfield_rows(sub: pd.DataFrame, col: str, max_items=3) -> List[List[str]]:
    rows=[]
    total=len(sub)
    if total == 0: return rows
    for val, g in sub.groupby(col):
        if val in {"", "NAN"}: val="-"
        m=core_metrics(g)
        rows.append([val, m["run_pct"], m["pass_pct"], pct(len(g), total)])
    rows.sort(key=lambda r: (-len(sub[sub[col]==r[0]]), r[0]))
    return rows[:max_items]

def explosive_summary(df: pd.DataFrame) -> Tuple[List[List[str]], List[List[str]], List[List[str]]]:
    run_exp=df[df["IS_RUN"] & df["IS_EXPLOSIVE"]]
    pass_exp=df[df["IS_PASS"] & df["IS_EXPLOSIVE"]]
    all_exp=df[df["IS_EXPLOSIVE"]]
    total=len(df)
    rows=[
        ["Run", str(len(run_exp)), pct(len(run_exp), total), ypp(run_exp["YARDS"].sum(), len(run_exp))],
        ["Pass", str(len(pass_exp)), pct(len(pass_exp), total), ypp(pass_exp["YARDS"].sum(), len(pass_exp))],
        ["All Explosives", str(len(all_exp)), pct(len(all_exp), total), ypp(all_exp["YARDS"].sum(), len(all_exp))],
    ]
    run_details=[]
    for _,r in run_exp.sort_values("YARDS", ascending=False).iterrows():
        run_details.append([r["PLAYCALL_RUN"], r["DND"], str(int(r["YARDS"]))])
    pass_details=[]
    for _,r in pass_exp.sort_values("YARDS", ascending=False).iterrows():
        pass_details.append([r["PLAYCALL_PASS"], r["DND"], str(int(r["YARDS"]))])
    return rows, run_details, pass_details

def run_direction(df: pd.DataFrame, n=8) -> List[List[str]]:
    run=df[df["IS_RUN"]].copy()
    top=run["OFF_PLAY"].value_counts().head(n).index.tolist()
    rows=[]
    for play in top:
        sub=run[run["OFF_PLAY"]==play]
        total=len(sub)
        combos=[
            ("L/Str L", (sub["PLAY_DIR"]=="L") & (sub["OFF_STR"]=="L")),
            ("L/Str R", (sub["PLAY_DIR"]=="L") & (sub["OFF_STR"]=="R")),
            ("R/Str R", (sub["PLAY_DIR"]=="R") & (sub["OFF_STR"]=="R")),
            ("R/Str L", (sub["PLAY_DIR"]=="R") & (sub["OFF_STR"]=="L")),
        ]
        vals=[]
        for label, mask in combos:
            c=int(mask.sum())
            vals.append(f"{c} / {pct(c,total)}")
        to_str=int(((sub["PLAY_DIR"]=="L") & (sub["OFF_STR"]=="L") | ((sub["PLAY_DIR"]=="R") & (sub["OFF_STR"]=="R"))).sum())
        rows.append([play] + vals + [pct(to_str,total)])
    return rows

def touch_analysis(df: pd.DataFrame, n=8) -> List[List[str]]:
    touch=df[df["PLAYER"].notna() & (df["PLAYER"]!="-")].copy()
    rows=[]
    total=len(touch)
    if total==0:
        return []
    for player, sub in touch.groupby("PLAYER"):
        plays=len(sub)
        top_play = sub["OFF_PLAY"].value_counts().index[0] if len(sub) else "-"
        rows.append([player, str(plays), pct(plays,total), ypp(sub["YARDS"].sum(), plays), pct(sub["EFF_BOOL"].sum(), plays), str(int(sub["IS_EXPLOSIVE"].sum())), top_play])
    rows.sort(key=lambda r: (-int(r[1]), r[0]))
    return rows[:n]

def situation_detail(df: pd.DataFrame, situation: str) -> Dict[str, Any]:
    sub=df[df["DND"]==situation]
    m=core_metrics(sub)
    def pct_eff_ypp(col, n=3):
        rows=[]
        if sub.empty: return rows
        for val,g in sub.groupby(col):
            rows.append([val, pct(len(g), len(sub)), pct(g["EFF_BOOL"].sum(), len(g)), ypp(g["YARDS"].sum(), len(g))])
        rows.sort(key=lambda r: (-pct_value(r[1]) if r[1]!="N/A" else 0, r[0]))
        return rows[:n]
    overall=[]
    if not sub.empty:
        for combo,g in sub.groupby(["OFF_FORM","FORM_TAG","BACKFIELD","MOTION","OFF_PLAY","PLAY_TAG"]):
            label=" / ".join(map(str, combo))
            overall.append([label, str(len(g)), pct(len(g), len(sub))])
        overall.sort(key=lambda r: (-int(r[1]), r[0]))
    return {
        "situation": situation,
        "metrics": m,
        "plays": pct_eff_ypp("OFF_PLAY"),
        "formations": pct_eff_ypp("OFF_FORM"),
        "personnel": pct_eff_ypp("PERSONNEL"),
        "players": pct_eff_ypp("PLAYER"),
        "field_routes": pct_eff_ypp("FIELD_ROUTE"),
        "boundary_routes": pct_eff_ypp("BOUNDARY_ROUTE"),
        "overall": overall[:5],
    }

def prediction_rows(df: pd.DataFrame, min_sample=3, max_rows=9) -> List[List[str]]:
    # Granular slice: Personnel / Formation / Field Zone / DND
    rows=[]
    cols=["PERSONNEL","OFF_FORM","FIELD_ZONE","DND"]
    for keys, sub in df.groupby(cols):
        n=len(sub)
        if n < min_sample: 
            continue
        m=core_metrics(sub)
        run_share=m["run"]/n if n else 0
        pass_share=m["pass"]/n if n else 0
        lean="Run" if run_share >= pass_share else "Pass"
        conf=max(run_share, pass_share)
        top_play=sub["OFF_PLAY"].value_counts().index[0] if n else "-"
        play_conf=(sub["OFF_PLAY"]==top_play).sum()/n if n else 0
        slice_label=f"{keys[0]} / {keys[1]} / {keys[2]} / {keys[3]}"
        rows.append([slice_label, str(n), lean, f"{conf*100:.0f}%", top_play, f"{play_conf*100:.0f}%"])
    rows.sort(key=lambda r: (-int(r[1]), -pct_value(r[3]), r[0]))
    return rows[:max_rows]

def predictability_rows(df: pd.DataFrame, group_cols: List[str], label_join=" / ", min_sample=3, max_rows=10) -> List[List[str]]:
    rows=[]
    for keys, sub in df.groupby(group_cols):
        if len(sub) < min_sample:
            continue
        m=core_metrics(sub)
        lean="Run" if m["run"] >= m["pass"] else "Pass"
        conf=max(m["run"], m["pass"])/len(sub)
        score=predictability_score(conf)
        label=label_join.join(keys) if isinstance(keys, tuple) else str(keys)
        rows.append([label, str(len(sub)), lean, f"{conf*100:.0f}%", score])
    rows.sort(key=lambda r: (-pct_value(r[3]), -int(r[1]), r[0]))
    return rows[:max_rows]

def predictability_score(share: float) -> str:
    val=share*100
    if val >= 90: return "Extremely Predictable"
    if val >= 80: return "Highly Predictable"
    if val >= 70: return "Strong Tendency"
    if val >= 60: return "Slight Tendency"
    return "Balanced"

def screen_alerts(df: pd.DataFrame) -> Tuple[str, List[List[str]]]:
    mask = df["OFF_PLAY"].str.contains("SCREEN", case=False, na=False) | df["PLAY_TAG"].str.contains("SCREEN", case=False, na=False)
    screens=df[mask]
    if screens.empty:
        return "No OFF PLAY or PLAY TAG entries contained “screen” in the ODK = O sample.", [["0", "No screen-tagged plays in upload"]]
    note=f"{len(screens)} screen-tagged plays found. Most common: {screens['OFF_PLAY'].value_counts().index[0]}."
    return note, [[str(len(screens)), f"{pct(len(screens), len(df))} of offense; top situation {screens['DND'].value_counts().index[0]}"]]

def defensive_alerts(df: pd.DataFrame) -> List[List[str]]:
    alerts=[]
    # Most predictable slice
    preds=prediction_rows(df, min_sample=3, max_rows=3)
    for r in preds[:2]:
        alerts.append([r[0], f"{r[2]} lean {r[3]}; top play {r[4]} ({r[5]})."])
    # Top formation/personnel
    if len(df):
        top_p=df["PERSONNEL"].value_counts().index[0]
        sub=df[df["PERSONNEL"]==top_p]; m=core_metrics(sub)
        alerts.append([f"Most-used personnel: {top_p}", f"{m['run_pct']} run across {m['total']} snaps."])
        top_f=df["OFF_FORM"].value_counts().index[0]
        sub=df[df["OFF_FORM"]==top_f]; m=core_metrics(sub)
        alerts.append([f"Most-used formation: {top_f}", f"{m['run_pct']} run across {m['total']} snaps."])
    return alerts[:3]

def executive_identity(df: pd.DataFrame) -> List[List[str]]:
    m=core_metrics(df)
    top_p=df["PERSONNEL"].value_counts().index[0] if len(df) else "-"
    top_p_sub=df[df["PERSONNEL"]==top_p]
    top_f=df["OFF_FORM"].value_counts().index[0] if len(df) else "-"
    top_f_sub=df[df["OFF_FORM"]==top_f]
    top_play=df["OFF_PLAY"].value_counts().index[0] if len(df) else "-"
    top_play_count=int((df["OFF_PLAY"]==top_play).sum()) if len(df) else 0
    rows=[
        ["Base call", f"{m['run_pct']} Run", "Load/run fit if run rate is 60%+; force throws to earn explosives." if m["run"]>=m["pass"] else "Pass-first profile; affect the QB and challenge route timing."],
        ["Top personnel", f"{top_p} • {len(top_p_sub)} snaps", f"{core_metrics(top_p_sub)['run_pct']} run in most-used grouping."],
        ["Top formation", f"{top_f} • {len(top_f_sub)} snaps", f"{core_metrics(top_f_sub)['run_pct']} run / {core_metrics(top_f_sub)['pass_pct']} pass."],
        ["Most called play", f"{top_play} • {top_play_count}x", f"{pct(top_play_count, len(df))} of all O snaps."],
    ]
    return rows

# -----------------------------
# POWERPOINT HELPERS
# -----------------------------

def iter_tables(slide):
    return [s for s in slide.shapes if getattr(s, "has_table", False)]

def set_shape_text(shape, text: str):
    """
    Replace shape text while preserving the formatting already present in the template.

    Important:
    Do NOT call text_frame.clear() for template titles. Clearing can reset title
    placeholder formatting and cause duplicated-section titles to turn white.
    This function keeps the first existing run and only changes its characters.
    """
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    text = str(text)

    # Preserve first paragraph/run formatting whenever possible.
    if tf.paragraphs:
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = text
            for extra in p.runs[1:]:
                extra.text = ""
        else:
            run = p.add_run()
            run.text = text
        # Clear additional paragraphs without removing formatting from the first.
        for extra_p in tf.paragraphs[1:]:
            for run in extra_p.runs:
                run.text = ""
        try:
            p.alignment = p.alignment
        except Exception:
            pass
    else:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text

def set_cell_text(cell, text: Any):
    cell.text = ""
    p=cell.text_frame.paragraphs[0]
    run=p.add_run()
    run.text=str(text)
    try:
        run.font.size = Pt(8)
    except Exception:
        pass

def fill_table(table_shape, rows: List[List[Any]], headers: Optional[List[str]]=None):
    tbl=table_shape.table
    nrows=len(tbl.rows)
    ncols=len(tbl.columns)
    if headers:
        for c,h in enumerate(headers[:ncols]):
            set_cell_text(tbl.cell(0,c), h)
    # clear body
    for r in range(1,nrows):
        for c in range(ncols):
            set_cell_text(tbl.cell(r,c), "")
    for r_idx,row in enumerate(rows[:max(0,nrows-1)], start=1):
        for c_idx,val in enumerate(row[:ncols]):
            set_cell_text(tbl.cell(r_idx,c_idx), val)

def fill_single_row_table(table_shape, values: List[Any]):
    tbl=table_shape.table
    row_idx=1 if len(tbl.rows)>1 else 0
    for c,val in enumerate(values[:len(tbl.columns)]):
        set_cell_text(tbl.cell(row_idx,c), val)

def duplicate_slide(prs: Presentation, slide_index: int, insert_after: Optional[int]=None):
    """Duplicate slide by index. Uses python-pptx private XML APIs."""
    source = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]
    dest = prs.slides.add_slide(blank_layout)

    # remove default shapes
    for shp in list(dest.shapes):
        dest.shapes._spTree.remove(shp.element)

    for shp in source.shapes:
        newel = copy.deepcopy(shp.element)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            dest.part.rels._add_relationship(rel.reltype, rel._target, rel.rId, rel.is_external)
        except Exception:
            pass

    if insert_after is not None:
        sldIdLst = prs.slides._sldIdLst
        sldId = sldIdLst[-1]
        sldIdLst.remove(sldId)
        sldIdLst.insert(insert_after+1, sldId)
        return prs.slides[insert_after+1]
    return dest

def populate_report(df: pd.DataFrame, template: str, output: str, opponent: str = "Opponent", games_label: str = "ODK = O only"):
    prs=Presentation(template)

    # Slide 1/2 Executive
    for idx in [0,1]:
        slide=prs.slides[idx]
        texts=[s for s in slide.shapes if getattr(s, "has_text_frame", False)]
        if idx == 0 and len(texts) >= 2:
            set_shape_text(texts[0], f"{opponent.upper()} OFFENSIVE\nTENDENCY REPORT")
            set_shape_text(texts[1], games_label)
        elif idx == 1 and len(texts) >= 2:
            set_shape_text(texts[0], "Executive Snapshot")
            set_shape_text(texts[1], "ODK = O")
        tables=iter_tables(slide)
        ident=executive_identity(df)
        metrics=core_metrics(df)
        if len(tables)>=2:
            fill_table(tables[0], ident)
            fill_single_row_table(tables[1], [
                f"{metrics['total']} | scouted snaps",
                f"{metrics['run_pct']} | {metrics['run']} runs",
                f"{metrics['pass_pct']} | {metrics['pass']} passes",
                metrics["eff_pct"],
                f"{metrics['explosive']} | {metrics['explosive_pct']}",
            ])

    # Slide 3 Part 1
    slide=prs.slides[2]
    tables=iter_tables(slide)
    fill_table(tables[0], run_pass_table(df, "PERSONNEL", max_rows=6))
    fill_table(tables[1], run_pass_table(df, "BACKFIELD", max_rows=8))
    fill_table(tables[2], dnd_table(df, max_rows=11))
    fill_table(tables[3], p10_table(df))
    fz=field_zone_table(df, max_rows=7)
    fill_table(tables[4], fz)
    fill_table(tables[5], fz)

    # Slide 4 Part 2 menu
    slide=prs.slides[3]
    tables=iter_tables(slide)
    fill_table(tables[0], formation_summary(df))
    fill_table(tables[1], combo_summary(df, 10))

    # Formation detail slides
    forms = [r[0] for r in formation_summary(df)]
    # duplicate slide 5 for additional forms, inserted after the current formation block
    form_start_idx=4
    for i in range(1, len(forms)):
        duplicate_slide(prs, form_start_idx, insert_after=form_start_idx+i-1)
    # Now populate formation slides
    for i,form in enumerate(forms):
        slide=prs.slides[form_start_idx+i]
        fd=formation_detail(df, form)
        shapes=[s for s in slide.shapes if getattr(s, "has_text_frame", False)]
        # Identify title-ish shapes by order
        for s in shapes:
            if "Formation Detail" in s.text or "DUO" in s.text:
                set_shape_text(s, f"Formation Detail — {form}")
                break
        tables=iter_tables(slide)
        if len(tables)>=6:
            fill_table(tables[0], fd["run_plays"])
            fill_table(tables[1], fd["field_routes"])
            fill_table(tables[2], fd["boundary_routes"])
            fill_table(tables[3], fd["motions"])
            fill_table(tables[4], fd["backfields"])
            m=fd["metrics"]
            fill_single_row_table(tables[5], [
                str(m["total"]),
                f"{m['run']} / {m['run_pct']}",
                f"{m['pass']} / {m['pass_pct']}",
                m["eff_pct"]
            ])
        # alert shapes by text
        for s in shapes:
            if s.text.strip().upper() == "DROP BACK" or "ALERT" in s.text:
                # hard to detect boxes; set the explicit alert text box only if it was the play label
                if s.text.strip().upper() == "DROP BACK":
                    set_shape_text(s, fd["alert"])
            elif "% of DUO" in s.text:
                set_shape_text(s, f"{fd['alert_pct']} of {form}")

    # Calculate dynamic indices after formation duplicates
    offset = max(0, len(forms)-1)
    s_exp=5+offset
    s_dir=6+offset
    s_touch=7+offset
    s_sit=8+offset
    s_pred=9+offset
    s_screen=10+offset
    s_score=11+offset
    s_plan=12+offset

    # Slide Explosives
    rows, run_det, pass_det=explosive_summary(df)
    tables=iter_tables(prs.slides[s_exp])
    fill_table(tables[0], rows)
    fill_table(tables[1], run_det)
    fill_table(tables[2], pass_det)

    # Slide Direction
    tables=iter_tables(prs.slides[s_dir])
    fill_table(tables[0], run_direction(df))

    # Slide Touch
    tables=iter_tables(prs.slides[s_touch])
    fill_table(tables[0], touch_analysis(df))

    # Situations: duplicate slide 9 for every DND present
    situations=[s for s in DND_ORDER if len(df[df["DND"]==s])]
    for i in range(1, len(situations)):
        duplicate_slide(prs, s_sit, insert_after=s_sit+i-1)
    for i,situation in enumerate(situations):
        slide=prs.slides[s_sit+i]
        sd=situation_detail(df, situation)
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and "Part 6" in shp.text:
                set_shape_text(shp, f"Part 6 — Situation Analysis: {situation}")
        tables=iter_tables(slide)
        if len(tables)>=8:
            fill_table(tables[0], sd["plays"])
            fill_table(tables[1], sd["formations"])
            fill_table(tables[2], sd["personnel"])
            fill_table(tables[3], sd["players"])
            fill_table(tables[4], sd["field_routes"])
            fill_table(tables[5], sd["boundary_routes"])
            fill_table(tables[6], sd["overall"])
            m=sd["metrics"]
            fill_single_row_table(tables[7], [str(m["total"]), f"{m['run_pct']} | {m['run']} calls", f"{m['pass_pct']} | {m['pass']} calls", f"{m['eff_pct']} efficient"])

    # Update indices after situations
    offset2=max(0, len(situations)-1)
    s_pred += offset2
    s_screen += offset2
    s_score += offset2
    s_plan += offset2

    # What's Next
    tables=iter_tables(prs.slides[s_pred])
    fill_table(tables[0], prediction_rows(df, min_sample=3, max_rows=9))

    # Screens
    note, screen_rows = screen_alerts(df)
    tables=iter_tables(prs.slides[s_screen])
    if tables:
        fill_table(tables[0], [[note]])
    if len(tables)>1:
        fill_table(tables[1], screen_rows)

    # Predictability
    tables=iter_tables(prs.slides[s_score])
    if len(tables)>=3:
        fill_table(tables[0], predictability_rows(df, ["DND"], min_sample=3, max_rows=8))
        fill_table(tables[1], predictability_rows(df, ["OFF_FORM"], min_sample=3, max_rows=6))
        fill_table(tables[2], predictability_rows(df, ["DND","OFF_FORM"], min_sample=3, max_rows=10))

    # Defensive Plan
    tables=iter_tables(prs.slides[s_plan])
    if tables:
        fill_table(tables[0], defensive_alerts(df))

    prs.save(output)

def generate(excel_paths: List[str], template: str, output: str, opponent: str, odk: str="O") -> Dict[str, Any]:
    raw=read_files(excel_paths)
    issues=validate(raw)
    df=add_features(raw, odk=odk)
    if df.empty:
        raise ValueError("No valid ODK=O Run/Pass plays found after filtering.")
    games_label=f"ODK = {odk.upper()} only • " + " + ".join([Path(p).stem for p in excel_paths])
    populate_report(df, template, output, opponent=opponent, games_label=games_label)
    m=core_metrics(df)
    return {"output": output, "issues": issues, "plays": m["total"], "run_pct": m["run_pct"], "pass_pct": m["pass_pct"], "explosives": m["explosive"]}

def main():
    parser=argparse.ArgumentParser(description="Generate DEF offensive tendency report from Hudl exports.")
    parser.add_argument("--template", required=True, help="Path to MASTER Offensive Breakdown Template.pptx")
    parser.add_argument("--output", required=True, help="Output PowerPoint path")
    parser.add_argument("--opponent", default="Opponent", help="Opponent/report title")
    parser.add_argument("--odk", default="O", help="ODK filter value; default O")
    parser.add_argument("--excel", nargs="+", required=True, help="One or more Excel/CSV files")
    args=parser.parse_args()
    result=generate(args.excel, args.template, args.output, args.opponent, args.odk)
    print("Report generated:", result)

if __name__ == "__main__":
    main()
